import os
import tempfile
import logging
import mimetypes
from pathlib import Path

class TextExtractor:
    """
    Class to extract text from different file types
    """
    
    def __init__(self):
        """Initialize TextExtractor"""
        self.logger = logging.getLogger(__name__)
    
    def extract_text(self, file_data, filename):
        """
        Extract text from file data based on file type
        
        Args:
            file_data (bytes): Binary file data
            filename (str): Original filename with extension
            
        Returns:
            str: Extracted text content
        """
        if not file_data:
            self.logger.warning("No file data provided for text extraction")
            return ""
            
        # Get file extension and mime type
        ext = os.path.splitext(filename)[1].lower()
        mime_type, _ = mimetypes.guess_type(filename)
        
        self.logger.info(f"Extracting text from {filename} ({mime_type})")
        
        # Handle different file types
        try:
            # Text files
            if mime_type and mime_type.startswith('text/'):
                return self._extract_from_text_file(file_data)
                
            # PDF files
            elif mime_type == 'application/pdf' or ext == '.pdf':
                return self._extract_from_pdf(file_data)
                
            # Microsoft Office documents
            elif mime_type and ('officedocument' in mime_type or 
                               mime_type in ['application/msword', 'application/vnd.ms-excel', 'application/vnd.ms-powerpoint']):
                return self._extract_from_office_document(file_data, mime_type, ext)
                
            # HTML files
            elif mime_type in ['text/html', 'application/xhtml+xml'] or ext in ['.html', '.htm', '.xhtml']:
                return self._extract_from_html(file_data)
                
            # CSV, JSON, XML
            elif ext in ['.csv', '.json', '.xml']:
                return self._extract_from_text_file(file_data)
                
            # Image files (OCR would be ideal but requires additional dependencies)
            elif mime_type and mime_type.startswith('image/'):
                self.logger.warning(f"Image file detected but OCR is not implemented: {filename}")
                return f"Image file: {filename}"
                
            # Fallback - try as UTF-8 text
            else:
                self.logger.warning(f"Unsupported file type: {mime_type}, trying as text")
                return self._extract_as_text(file_data)
                
        except Exception as e:
            self.logger.error(f"Error extracting text from {filename}: {str(e)}")
            return f"Error extracting text: {str(e)}"
    
    def _extract_from_text_file(self, file_data):
        """Extract text from a text file"""
        try:
            return file_data.decode('utf-8', errors='replace')
        except UnicodeDecodeError:
            # Try other common encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    return file_data.decode(encoding, errors='replace')
                except UnicodeDecodeError:
                    continue
            # Fallback
            return file_data.decode('utf-8', errors='replace')
    
    def _extract_from_pdf(self, file_data):
        """Extract text from a PDF file"""
        try:
            # Create a temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
                temp_path = temp_file.name
                temp_file.write(file_data)
            
            try:
                # Try using PyPDF2 if available
                import PyPDF2
                text = ""
                with open(temp_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    for page_num in range(len(reader.pages)):
                        text += reader.pages[page_num].extract_text() + "\n"
                return text
            except ImportError:
                # Try using pdfplumber if available
                try:
                    import pdfplumber
                    text = ""
                    with pdfplumber.open(temp_path) as pdf:
                        for page in pdf.pages:
                            text += page.extract_text() or "" + "\n"
                    return text
                except ImportError:
                    # If neither library is available
                    self.logger.warning("PDF extraction libraries not available")
                    return "PDF file (text extraction libraries not available)"
        finally:
            # Clean up the temporary file
            if 'temp_path' in locals() and os.path.exists(temp_path):
                os.unlink(temp_path)
    
    def _extract_from_office_document(self, file_data, mime_type, ext):
        """Extract text from Microsoft Office documents"""
        try:
            # Create a temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as temp_file:
                temp_path = temp_file.name
                temp_file.write(file_data)
            
            try:
                # Try using different libraries based on file type
                
                # Word documents
                if '.doc' in ext or 'word' in mime_type:
                    try:
                        # Try docx2txt first (for .docx)
                        import docx2txt
                        return docx2txt.process(temp_path)
                    except (ImportError, Exception):
                        # Try python-docx
                        try:
                            import docx
                            doc = docx.Document(temp_path)
                            return "\n".join([para.text for para in doc.paragraphs])
                        except ImportError:
                            self.logger.warning("Word document libraries not available")
                            return "Word document (text extraction libraries not available)"
                
                # Excel documents
                elif '.xls' in ext or 'excel' in mime_type:
                    try:
                        import pandas as pd
                        df = pd.read_excel(temp_path)
                        return df.to_string()
                    except ImportError:
                        self.logger.warning("Excel libraries not available")
                        return "Excel document (text extraction libraries not available)"
                
                # PowerPoint documents
                elif '.ppt' in ext or 'powerpoint' in mime_type:
                    try:
                        import pptx
                        presentation = pptx.Presentation(temp_path)
                        text = ""
                        for slide in presentation.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text += shape.text + "\n"
                        return text
                    except ImportError:
                        self.logger.warning("PowerPoint libraries not available")
                        return "PowerPoint document (text extraction libraries not available)"
                
                # Default case
                return "Microsoft Office document (type not specifically supported)"
                
            except Exception as e:
                self.logger.error(f"Error extracting from Office document: {str(e)}")
                return f"Error extracting from Office document: {str(e)}"
        finally:
            # Clean up the temporary file
            if 'temp_path' in locals() and os.path.exists(temp_path):
                os.unlink(temp_path)
    
    def _extract_from_html(self, file_data):
        """Extract text from HTML content"""
        try:
            # Try using BeautifulSoup if available
            from bs4 import BeautifulSoup
            html_content = file_data.decode('utf-8', errors='replace')
            soup = BeautifulSoup(html_content, 'html.parser')
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.extract()
            # Get text
            text = soup.get_text()
            # Break into lines and remove leading and trailing space
            lines = (line.strip() for line in text.splitlines())
            # Break multi-headlines into a line each
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            # Drop blank lines
            text = '\n'.join(chunk for chunk in chunks if chunk)
            return text
        except ImportError:
            # If BeautifulSoup is not available, try a simple approach
            html_content = file_data.decode('utf-8', errors='replace')
            # Remove HTML tags (very basic approach)
            import re
            text = re.sub(r'<[^>]+>', ' ', html_content)
            return text
    
    def _extract_as_text(self, file_data):
        """Extract as plain text, last resort"""
        try:
            return file_data.decode('utf-8', errors='replace')
        except:
            return "Binary file (text extraction failed)"